import os
import re
import time
import textwrap
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from google import genai
from google.genai import types
from PIL import Image

google_client = genai.Client(api_key=os.getenv("GOOGLE_GEMINI_API_KEY"))

def get_subtopics(topic):
    prompt = (
        f"List exactly 6 concise subtopics for the main topic '{topic}'. "
        "Respond with only the subtopics as a numbered list—without any introductions, explanations, or additional text."
    )
    response = google_client.models.generate_content(
        model="gemini-2.0-flash",
        contents=[prompt],
        config=types.GenerateContentConfig(response_modalities=['TEXT'])
    )
    subtopics = [
        re.sub(r"^\d+\.\s*", "", line.strip())
        for line in response.candidates[0].content.parts[0].text.strip().split("\n")
        if line.strip()
    ][:6]
    return subtopics

def generate_description(subtopic, topic):
    prompt = f"Write a concise one-paragraph description of 1100 characters for the subtopic '{subtopic}' under '{topic}', without introductory phrases."
    response = google_client.models.generate_content(
        model="gemini-2.0-flash",
        contents=[prompt],
        config=types.GenerateContentConfig(response_modalities=['TEXT'])
    )
    return response.candidates[0].content.parts[0].text.strip()

def generate_image(subtopic, topic):
    image_prompt = f"Create a hyper-realistic and intricately detailed 3D rendered image that artistically represents the concept of '{subtopic}' within the context of '{topic}'. The image should depict a dynamic scene filled with realistic elements, natural lighting, and rich textures that emphasize the emotional and practical significance of '{subtopic}'—using symbolic visual cues rather than any text or labels."
    generated_image = None
    attempts = 0
    max_attempts = 5
    while attempts < max_attempts and generated_image is None:
        try:
            image_response = google_client.models.generate_content(
                model="gemini-2.0-flash-exp-image-generation",
                contents=[image_prompt],
                config=types.GenerateContentConfig(response_modalities=['TEXT', 'IMAGE'])
            )
            for part in image_response.candidates[0].content.parts:
                if part.inline_data is not None:
                    generated_image = Image.open(BytesIO(part.inline_data.data))
                    break
        except genai.errors.ServerError as e:
            if "The model is overloaded" in str(e):
                print(f"Attempt {attempts+1} failed: The model is overloaded. Please try again later. Retrying in 5 seconds...")
            else:
                print(f"Attempt {attempts+1} failed: {str(e)}. Retrying in 5 seconds...")
            time.sleep(5)
            attempts += 1
    return generated_image

def build_presentation(slides_info, topic):
    prs = Presentation()
    slide_width_inch = prs.slide_width / 914400.0
    slide_height_inch = prs.slide_height / 914400.0
    
    for idx, slide_data in enumerate(slides_info, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 248, 255)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        title_box.text_frame.text = slide_data["subtopic"]
        for paragraph in title_box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(28)
            paragraph.alignment = PP_ALIGN.CENTER
        
        desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(2.8))  
        desc_frame = desc_box.text_frame
        desc_frame.paragraphs[0].text = textwrap.fill(slide_data["description"], width=95)
        desc_frame.paragraphs[0].font.size = Pt(16)
        desc_frame.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
        
        if slide_data["image"]:
            orig_width, orig_height = slide_data["image"].width, slide_data["image"].height
            aspect_ratio = orig_width / orig_height
            final_width = min(6.5, 3.4 * aspect_ratio)
            final_height = final_width / aspect_ratio
            image_left = (slide_width_inch - final_width) / 2.0
            image_top = 3.6 + (slide_height_inch - 3.6 - final_height) / 2.0
            image_stream = BytesIO()
            slide_data["image"].save(image_stream, format='PNG')
            image_stream.seek(0)
            slide.shapes.add_picture(image_stream, Inches(image_left), Inches(image_top),
                                       width=Inches(final_width), height=Inches(final_height))
    
    ppt_filename = f"{topic}.pptx"
    prs.save(ppt_filename)
    print(f"Presentation saved as {ppt_filename}")
    return ppt_filename

def main():
    topic_input = input("Enter the main topic for the presentation: ").strip()
    topic = topic_input if topic_input else "Sustainable Energy"
    
    subtopics = get_subtopics(topic)
    slides_info = []
    for subtopic in subtopics:
        description = generate_description(subtopic, topic)
        image = generate_image(subtopic, topic)
        slides_info.append({
            "subtopic": subtopic,
            "description": description,
            "image": image
        })
    
    ppt_filename = build_presentation(slides_info, topic)
    os.startfile(ppt_filename)

if __name__ == "__main__":
    main()
